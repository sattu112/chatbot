import os
import fitz  # PyMuPDF
from PIL import Image
import pytesseract
import pandas as pd
from docx import Document

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower().strip()
    text = ""
    try:
        if ext == '.pdf':
            text = extract_text_from_pdf(file_path)
        elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            text = extract_text_from_image(file_path)
        elif ext in ['.xls', '.xlsx']:
            text = extract_text_from_excel(file_path)
        elif ext in ['.doc', '.docx']:
            text = extract_text_from_word(file_path)
        elif ext in ['.ppt', '.pptx']:
            text = extract_text_from_ppt(file_path)
        elif ext in ['.txt', '.csv']:
            text = extract_text_from_text(file_path)
        else:
            text = f"Unsupported file format: {ext}"
    except Exception as e:
        text = f"Error extracting text: {str(e)}"
    return text

def extract_text_from_pdf(file_path):
    text = ""
    pdf = fitz.open(file_path)
    for page in pdf:
        text += page.get_text() + "\n"
    pdf.close()
    return text

def extract_text_from_image(image_path):
    try:
        image = Image.open(image_path)
        return pytesseract.image_to_string(image)
    except Exception as e:
        return "OCR error: Please ensure Tesseract OCR is installed and in PATH."

def extract_text_from_excel(file_path):
    df = pd.read_excel(file_path, sheet_name=None)
    text = ""
    for sheet, data in df.items():
        text += f"Sheet: {sheet}\n" + data.to_string() + "\n"
    return text

def extract_text_from_word(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def extract_text_from_ppt(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

def extract_text_from_text(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()
